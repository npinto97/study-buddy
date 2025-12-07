import tempfile
import os
import base64
import uuid
import requests
import time
from typing import Union, List, Optional, Dict, Any
from abc import ABC, abstractmethod

from study_buddy.utils.logging_config import get_logger, LogContext, metrics

logger = get_logger("tools")

import pandas as pd
import fitz  # pip install PyMuPDF
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
from textblob import TextBlob
from docx import Document  # pip install python-docx

from pydantic import BaseModel, Field
from langchain_core.tools import Tool

# Compatibility: StructuredTool moved between langchain versions.
# Try to import it; if unavailable, provide a lightweight fallback
# that implements `from_function` used by this codebase.
try:
    from langchain.tools import StructuredTool
except Exception:
    class StructuredTool:
        """Lightweight fallback for langchain.tools.StructuredTool.

        Provides a `from_function` method that returns a `Tool` wrapping
        the provided function. If `args_schema` (a Pydantic model class)
        is provided, input will be validated before calling the function.
        This fallback is intentionally minimal and only intended to keep
        older/newer langchain installations compatible with the project.
        """

        @staticmethod
        def from_function(func, name: str, description: str, args_schema=None):
            def wrapper(*args, **kwargs):
                # If an args_schema (Pydantic model class) is provided, try to
                # validate inputs and call the wrapped function with the
                # validated dict. Support calling with kwargs, a single dict
                # positional, or positional arguments mapped to model fields.
                if args_schema is not None:
                    try:
                        # If kwargs supplied, validate directly
                        if kwargs:
                            validated = args_schema(**kwargs)
                            return func(**validated.dict())

                        # If a single dict positional was supplied
                        if len(args) == 1 and isinstance(args[0], dict):
                            validated = args_schema(**args[0])
                            return func(**validated.dict())

                        # Map positional args to model fields if possible
                        if len(args) >= 1 and hasattr(args_schema, "__fields__"):
                            fields = list(args_schema.__fields__.keys())
                            mapped = {fields[i]: args[i] for i in range(min(len(args), len(fields)))}
                            validated = args_schema(**mapped)
                            return func(**validated.dict())

                    except Exception:
                        # Re-raise validation errors to surface helpful messages
                        raise

                # Fallback: call the function directly
                return func(*args, **kwargs)

            return Tool(name=name, description=description, func=wrapper)
from langchain_community.document_loaders import PyPDFLoader, TextLoader, CSVLoader
from langchain_community.llms import Together
from langchain_community.tools.google_lens import GoogleLensQueryRun
from langchain_community.utilities.google_lens import GoogleLensAPIWrapper
from langchain_community.tools import YouTubeSearchTool, WikipediaQueryRun
from langchain_community.utilities import WikipediaAPIWrapper
from langchain_community.tools.google_books import GoogleBooksQueryRun
from langchain_community.utilities.google_books import GoogleBooksAPIWrapper
from langchain_community.tools.google_scholar import GoogleScholarQueryRun
from langchain_community.utilities.google_scholar import GoogleScholarAPIWrapper
# Tavily may not be installed in all environments; provide a lightweight
# fallback that reports the feature is unavailable. This prevents import
# errors during app startup and allows the rest of the tools to load.
try:
    from langchain_tavily import TavilySearch
except Exception:
    class TavilySearch:
        def __init__(self, *args, **kwargs):
            self.available = False

        def run(self, *args, **kwargs):
            return (
                "TavilySearch is not installed in this environment. "
                "Install the 'langchain-tavily' package or set up an alternative web search provider."
            )
# from langchain_community.tools.arxiv.tool import ArxivQueryRun
try:
    from langchain.text_splitter import CharacterTextSplitter
except ImportError:
    from langchain_text_splitters import CharacterTextSplitter

from e2b_code_interpreter import Sandbox
from elevenlabs.client import ElevenLabs
from elevenlabs import save
import assemblyai as aai
from together import Together as TogetherClient

from study_buddy.vectorstore_pipeline.vector_store_builder import get_vector_store
from study_buddy.config import FAISS_INDEX_DIR


# =============================================================================
# Configuration and Constants
# =============================================================================

class Config:
    """Centralized configuration management."""
    GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
    TOGETHER_API_KEY = os.getenv("TOGETHER_API_KEY")
    ELEVEN_API_KEY = os.getenv("ELEVEN_API_KEY")
    ASSEMBLYAI_API_KEY = os.getenv("ASSEMBLYAI_API_KEY")
    SPOTIFY_CLIENT_ID = os.getenv("SPOTIFY_CLIENT_ID")
    SPOTIFY_CLIENT_SECRET = os.getenv("SPOTIFY_CLIENT_SECRET")
    E2B_API_KEY = os.getenv("E2B_API_KEY")
    SERP_API_KEY = os.getenv("SERP_API_KEY")
    TAVILY_API_KEY = os.getenv("TAVILY_API_KEY")
    GOOGLE_LENS_API_KEY = os.getenv("GOOGLE_LENS_API_KEY")
    IMGBB_API_KEY = os.getenv("IMGBB_API_KEY")  # For temporary image hosting
    
    # Default models and settings
    DEFAULT_LLM_MODEL = "mistralai/Mistral-7B-Instruct-v0.2"
    DEFAULT_TTS_MODEL = "eleven_multilingual_v2"
    DEFAULT_VOICE_ID = "paraDwhkbSkX4FhBkAzc"
    
    @classmethod
    def validate_key(cls, key_name: str) -> str:
        """Validate that required API key exists."""
        key = getattr(cls, key_name, None)
        if not key:
            raise ValueError(f"{key_name} must be set in environment variables")
        return key


# =============================================================================
# Base Classes and Utilities
# =============================================================================

class BaseWrapper(ABC):
    """Base class for all tool wrappers."""
    
    def __init__(self, **kwargs):
        self.validate_dependencies()
    
    @abstractmethod
    def validate_dependencies(self):
        """Validate required dependencies and API keys."""
        pass


class FileProcessor:
    """Utility class for file operations."""
    
    @staticmethod
    def get_file_extension(file_path: str) -> str:
        """Get normalized file extension."""
        return os.path.splitext(file_path)[1].lower()
    
    @staticmethod
    def create_temp_file(suffix: str = "") -> str:
        """Create temporary file and return path."""
        return tempfile.NamedTemporaryFile(delete=False, suffix=suffix).name
    
    @staticmethod
    def download_file(url: str, suffix: str = "") -> str:
        """Download file from URL to temporary location."""
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        temp_path = FileProcessor.create_temp_file(suffix)
        with open(temp_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
        return temp_path
    
    @staticmethod
    def upload_image_to_imgbb(image_path: str, api_key: Optional[str] = None) -> Optional[str]:
        """
        Upload a local image to imgbb.com and return the public URL.
        
        Args:
            image_path: Path to the local image file
            api_key: imgbb API key (optional, will use Config.IMGBB_API_KEY if not provided)
            
        Returns:
            Public URL of the uploaded image, or None if upload fails
        """
        if api_key is None:
            api_key = Config.IMGBB_API_KEY
        
        if not api_key:
            logger.warning("IMGBB_API_KEY not set, cannot upload image")
            return None
        
        if not os.path.exists(image_path):
            logger.error(f"Image file not found: {image_path}")
            return None
        
        try:
            # Read and encode image to base64
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode("utf-8")
            
            # Upload to imgbb
            url = "https://api.imgbb.com/1/upload"
            payload = {
                "key": api_key,
                "image": image_data,
                "expiration": 600  # 10 minutes expiration
            }
            
            logger.info(f"Uploading image to imgbb: {os.path.basename(image_path)}")
            response = requests.post(url, data=payload, timeout=30)
            response.raise_for_status()
            
            result = response.json()
            if result.get("success"):
                public_url = result["data"]["url"]
                logger.info(f"Image uploaded successfully: {public_url}")
                return public_url
            else:
                logger.error(f"imgbb upload failed: {result}")
                return None
                
        except Exception as e:
            logger.error(f"Failed to upload image to imgbb: {e}", exc_info=True)
            return None


class LLMFactory:
    """Factory for creating LLM instances."""
    
    @staticmethod
    def create_together_llm(
        model: str = Config.DEFAULT_LLM_MODEL,
        temperature: float = 0.2,
        max_tokens: int = 1024
    ):
        """Create Together LLM instance."""
        return Together(
            model=model,
            temperature=temperature,
            max_tokens=max_tokens,
            together_api_key=Config.validate_key("TOGETHER_API_KEY")
        )
        

def resolve_file_path(file_path: str) -> str:
    """
    Risolve e sanitizza un percorso di file, gestendo percorsi assoluti,
    relativi e stringhe potenzialmente malformate.
    """
    if not isinstance(file_path, str):
        logger.error(f"Invalid file path type: {type(file_path)}. Must be a string.")
        return ""

    # 1. Rimuovi eventuali apici o virgolette che l'LLM potrebbe aver aggiunto
    clean_path = file_path.strip().strip("'\"")

    # 2. Normalizza i separatori di percorso per il sistema operativo corrente
    if os.name == 'nt':  # Windows
        clean_path = clean_path.replace('/', '\\')
    else:  # Unix-like
        clean_path = clean_path.replace('\\', '/')
        
    # 3. Se il percorso è assoluto ed esiste, usalo
    if os.path.isabs(clean_path) and os.path.exists(clean_path):
        return os.path.normpath(clean_path)

    # 4. Altrimenti, cercalo nella cartella di upload predefinita
    upload_dir = os.path.join(os.getcwd(), "uploaded_files")
    # Usa os.path.basename per gestire in sicurezza sia percorsi relativi che nomi di file semplici
    potential_path = os.path.join(upload_dir, os.path.basename(clean_path))

    if os.path.exists(potential_path):
        return os.path.normpath(potential_path)

    # 5. Come ultima risorsa, restituisci il percorso pulito ma non verificato
    logger.warning(f"Could not find file in known paths: '{file_path}'. Returning normalized path.")
    return os.path.normpath(clean_path)

# =============================================================================
# Pydantic Models for Input Validation
# =============================================================================

class FilePathInput(BaseModel):
    """Base model for tools that require file paths."""
    file_path: str = Field(description="Path to the input file")


class QueryInput(BaseModel):
    """Base model for tools that require queries."""
    query: str = Field(description="Search or analysis query")


class CodeInterpreterInput(BaseModel):
    """Input schema for code interpreter."""
    code: str = Field(description="Python code to execute")


class DataVizInput(BaseModel):
    """Input schema for data visualization."""
    csv_path: str = Field(description="Absolute path to the CSV file")
    query: str = Field(description="Natural language question about the dataset")


class AudioInput(BaseModel):
    """Input schema for audio processing."""
    audio_input: Union[str, bytes] = Field(description="Audio file path, URL, or raw bytes")


class GoogleLensInput(BaseModel):
    """Input schema for Google Lens analysis."""
    query: str = Field(description="Path to an image file or a public image URL to analyze with Google Lens")


# =============================================================================
# Core Tool Implementations
# =============================================================================

class VectorStoreRetriever(BaseWrapper):
    
    def validate_dependencies(self):
        """Validate vector store availability."""
        if not os.path.exists(FAISS_INDEX_DIR):
            raise ValueError(f"FAISS index directory not found: {FAISS_INDEX_DIR}")
    
    def retrieve(self, query: str, k: int = 6, min_score: float = 0.1) -> tuple[str, list, list]:
        """
        Retrieve information with validation, enhanced logging, and quality filtering.
        
        Args:
            query: Search query text
            k: Number of results to retrieve (default: 4)
            min_score: Minimum similarity score threshold (default: 0.3)
            
        Returns:
            tuple of (formatted_results, doc_objects, file_paths)
        """
        with LogContext("rag_retrieval", logger) as log_ctx:
            logger.info(f"""
🔍 RAG Retrieval Process:
   Query: {query}
   Requested Results: {k}
   Min Score Threshold: {min_score}
""")
            
            try:
                # Initialize vector store with metrics
                init_start = time.time()
                vector_store = get_vector_store(FAISS_INDEX_DIR)
                init_time = time.time() - init_start
                
                logger.debug(f"""
📚 Vector Store Details:
   Location: {FAISS_INDEX_DIR}
   Type: {type(vector_store).__name__}
   Initialization Time: {init_time:.3f}s
""")
                
                # Validate embedding model
                embed_start = time.time()
                embedding_dim = vector_store.index.d
                test_embedding = vector_store.embeddings.embed_query("test")
                embed_time = time.time() - embed_start
                
                if len(test_embedding) != embedding_dim:
                    raise ValueError(f"Embedding dimension mismatch: {len(test_embedding)} vs {embedding_dim}")
                
                logger.debug(f"""
🧬 Embedding Validation:
   Dimension: {embedding_dim}
   Test Time: {embed_time:.3f}s
   Model: {type(vector_store.embeddings).__name__}
""")
                
                # Perform enhanced similarity search
                search_start = time.time()
                # Fetch more candidates for better filtering
                retrieved_docs = vector_store.similarity_search_with_score(
                    query, 
                    k=min(k * 2, 20)  # Fetch extra for quality filtering
                )
                search_time = time.time() - search_start
                
                # Quality-based filtering
                filtered_docs = []
                for doc, score in retrieved_docs:
                    if score >= min_score:
                        filtered_docs.append((doc, score))
                    else:
                        logger.debug(f"Filtered out document with low score: {score:.4f}")
                
                # Sort by score and take top k
                filtered_docs.sort(key=lambda x: x[1], reverse=True)
                final_docs = filtered_docs[:k]
                
                # HYBRID SEARCH: Detect contact-related queries and add keyword search
                contact_keywords = ['mail', 'email', 'e-mail', 'telefono', 'phone', 'contatto', 'contact']
                is_contact_query = any(keyword in query.lower() for keyword in contact_keywords)
                
                if is_contact_query:
                    logger.info("🔍 Contact query detected - performing hybrid search with keyword matching")
                    
                    # Extract person name from query (simple heuristic)
                    import re
                    # Look for common patterns like "professor X" or just names
                    name_patterns = [
                        r'(?:professor|prof\.?|dott\.?)\s+([a-zA-Z]+(?:\s+[a-zA-Z]+)*)',
                        r'(?:di|del)\s+([a-zA-Z]+(?:\s+[a-zA-Z]+)*)',
                    ]
                    person_name = None
                    for pattern in name_patterns:
                        match = re.search(pattern, query, re.IGNORECASE)
                        if match:
                            person_name = match.group(1).strip()
                            break
                    
                    # Fallback: take last capitalized word(s) as name
                    if not person_name:
                        words = query.split()
                        capitalized = [w for w in words if w and w[0].isupper()]
                        if capitalized:
                            person_name = ' '.join(capitalized[-2:]) if len(capitalized) >= 2 else capitalized[-1]
                    
                    logger.info(f"   Extracted name: {person_name}")
                    
                    # Search for documents containing both the name and contact patterns
                    docstore = vector_store.docstore._dict
                    keyword_matches = []
                    
                    for doc_id, doc in docstore.items():
                        content_lower = doc.page_content.lower()
                        
                        # Check if document contains the person's name
                        if person_name and person_name.lower() in content_lower:
                            # Check if it contains email or phone patterns
                            has_email = '@' in doc.page_content and any(domain in content_lower for domain in ['.it', '.com', '.edu'])
                            has_phone = bool(re.search(r'\d{3}[-\s]?\d{7}', doc.page_content))
                            
                            if has_email or has_phone:
                                # Calculate a synthetic score based on relevance
                                score = 1.5  # Higher than semantic search to prioritize
                                keyword_matches.append((doc, score))
                                logger.info(f"   ✓ Found contact info for {person_name} in {doc.metadata.get('file_path', 'Unknown')}")
                    
                    # Merge keyword matches with semantic results
                    if keyword_matches:
                        logger.success(f"   Found {len(keyword_matches)} documents via keyword search")
                        # Add keyword matches at the beginning (higher priority)
                        all_docs = keyword_matches + final_docs
                        # Remove duplicates (keep first occurrence)
                        seen_ids = set()
                        unique_docs = []
                        for doc, score in all_docs:
                            doc_id = id(doc)
                            if doc_id not in seen_ids:
                                seen_ids.add(doc_id)
                                unique_docs.append((doc, score))
                        final_docs = unique_docs[:k]
                        logger.info(f"   Hybrid results: {len(final_docs)} documents (keyword + semantic)")
                
                logger.info(f"""
🎯 Search Performance:
   Time: {search_time:.3f}s
   Raw Results: {len(retrieved_docs)}
   Filtered Results: {len(filtered_docs)}
   Final Results: {len(final_docs)}
   Avg Score: {sum(score for _, score in final_docs) / len(final_docs) if final_docs else 0:.4f}
""")
                
                # Process results with detailed metadata
                file_paths = []
                serialized_parts = []
                metrics = {"scores": [], "content_lengths": []}
                
                for i, (doc, score) in enumerate(final_docs, 1):
                    # Enhanced metadata processing
                    metadata = doc.metadata
                    path = metadata.get("file_path", "Unknown")
                    page = metadata.get("page", "N/A")
                    doc_type = metadata.get("type", "Unknown")
                    
                    # Use only filename instead of full absolute path for portability
                    if path and path != "Unknown":
                        filename = os.path.basename(path)  # Extract just the filename
                        normalized_path = os.path.normpath(path)
                        file_paths.append(normalized_path)
                    else:
                        filename = "Unknown"
                    
                    # Track metrics
                    metrics["scores"].append(score)
                    metrics["content_lengths"].append(len(doc.page_content))
                    
                    # Format detailed document info - using filename only
                    doc_info = f"""
Document {i}:
• Score: {score:.4f}
• Source: {filename}
• Page/Section: {page}
• Type: {doc_type}
• Length: {len(doc.page_content)} chars
• Content Preview:
{'-' * 40}
{doc.page_content[:200]}...
{'-' * 40}
"""
                    logger.info(doc_info)
                    
                    # Format for return - using filename only for portability
                    serialized_parts.append(
                        f"[Source {i}] {filename} (Page {page})\n"
                        f"Relevance: {score:.4f}\n"
                        f"{doc.page_content}\n"
                        f"{'-' * 80}"
                    )
                
                # Log final metrics
                logger.success(f"""
📊 Retrieval Metrics:
   Documents Retrieved: {len(final_docs)}
   Average Score: {sum(metrics['scores']) / len(metrics['scores']):.4f}
   Score Range: {min(metrics['scores']):.4f} - {max(metrics['scores']):.4f}
   Avg Content Length: {sum(metrics['content_lengths']) / len(metrics['content_lengths']):.0f}
   Total Process Time: {search_time + init_time + embed_time:.3f}s
""")
                
                combined_results = "\n\n".join(serialized_parts)
                return combined_results, [doc for doc, _ in final_docs], file_paths
                
            except Exception as e:
                logger.error(f"""
❌ Retrieval Failed:
   Error Type: {type(e).__name__}
   Error Message: {str(e)}
""", exc_info=True)
                
                return (
                    f"Error during retrieval: {type(e).__name__} - {str(e)}\n"
                    "Please try reformulating your query or contact support.",
                    [], []
                )
    def print_faithfulness_percent(query, k=4, min_score=0.3):
        """
        Esegue una query con VectorStoreRetriever e stampa la percentuale di faithfulness per ogni documento e la media.
        """
        retriever = VectorStoreRetriever()
        results, docs, paths = retriever.retrieve(query, k=k, min_score=min_score)
        print(f"Risultati faithfulness per la query: {query}\n")
        # Estrai gli score dal testo dei risultati (che contiene 'Relevance: <score>')
        import re
        pattern = r"Relevance: ([0-9.]+)"
        scores = [float(m) for m in re.findall(pattern, results)]
        if not scores:
            print("Nessun risultato trovato o impossibile estrarre gli score.")
            return
        for i, score in enumerate(scores, 1):
            print(f"Documento {i}: Faithfulness = {score*100:.2f}%")
        print(f"\nFaithfulness media: {sum(scores)/len(scores)*100:.2f}%\n")


class DocumentProcessor(BaseWrapper):
    """Unified document processing for various file types."""
    
    def validate_dependencies(self):
        """Validate document processing dependencies."""
        pass

    def extract_text(self, file_path: str) -> str:
        """
        Estrae il testo da vari tipi di documenti, delegando la risoluzione
        e la pulizia del percorso alla funzione `resolve_file_path`.
        Per documenti lunghi (>2000 caratteri), restituisce un preview strutturato
        invece del contenuto completo per evitare overflow di contesto.
        """
        logger.info(f"📄 extract_text called with raw input: {repr(file_path)}")
        
        resolved_path = resolve_file_path(file_path)
        logger.info(f"📄 Path resolved to: {repr(resolved_path)}")
        
        if not os.path.exists(resolved_path):
            # Check if it's a path issue (e.g. old project name 'univox' vs 'study-buddy')
            if 'univox' in resolved_path and 'study-buddy' in os.getcwd():
                logger.warning(f"Detected potential path mismatch: 'univox' in path but current dir is 'study-buddy'")
            
            error_message = (
                f"FILE NON TROVATO: Il sistema non riesce a leggere il file '{os.path.basename(resolved_path)}'.\n"
                f"Percorso cercato: {resolved_path}\n\n"
                "POSSIBILE CAUSA: Il file è stato indicizzato in precedenza (es. nel progetto 'univox') ma non è fisicamente presente nella cartella di lavoro attuale ('study-buddy').\n"
                "SOLUZIONE: Per analizzare questo file, devi caricarlo nuovamente tramite l'interfaccia ('Upload Document')."
            )
            logger.error(error_message)
            return error_message

        logger.info("✅ Path exists: True. Proceeding with text extraction.")
        
        try:
            ext = FileProcessor.get_file_extension(resolved_path)
            
            # Extract raw text based on file type
            if ext == ".pdf":
                raw_text = self._extract_from_pdf(resolved_path)
            elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
                raw_text = self._extract_from_image(resolved_path)
            elif ext == ".txt":
                raw_text = self._extract_from_text(resolved_path)
            elif ext in [".docx", ".doc"]:
                raw_text = self._extract_from_docx(resolved_path)
            elif ext in [".pptx", ".ppt"]:
                raw_text = self._extract_from_pptx(resolved_path)
            else:
                unsupported_error = f"Unsupported file format: '{ext}'. Cannot extract text from '{os.path.basename(resolved_path)}'."
                logger.warning(unsupported_error)
                return unsupported_error
            
            # For long documents, return a structured preview instead of full text
            # This prevents context overflow and helps the LLM understand the content better
            if len(raw_text) > 2000:
                logger.info(f"📝 Document is long ({len(raw_text)} chars), creating structured preview")
                
                # Create a structured preview with beginning, middle sample, and end
                lines = raw_text.split('\n')
                total_lines = len(lines)
                
                # Take ~25% from start, ~10% from middle, ~10% from end
                start_count = min(int(total_lines * 0.25), 30)
                middle_start = int(total_lines * 0.45)
                middle_count = min(int(total_lines * 0.10), 10)
                end_count = min(int(total_lines * 0.10), 10)
                
                start_text = '\n'.join(lines[:start_count])
                middle_text = '\n'.join(lines[middle_start:middle_start + middle_count])
                end_text = '\n'.join(lines[-end_count:]) if end_count > 0 else ""
                
                structured_preview = f"""DOCUMENTO ESTRATTO: {os.path.basename(resolved_path)}
Lunghezza totale: {len(raw_text)} caratteri, {total_lines} righe

=== INIZIO DOCUMENTO ===
{start_text}

=== SEZIONE CENTRALE (esempio) ===
{middle_text}

=== FINE DOCUMENTO ===
{end_text}

[Documento completo estratto ma mostrato in forma riassuntiva per evitare overflow]"""
                
                return structured_preview
            
            # For short documents, return full text
            return raw_text
                
        except Exception as e:
            extraction_error = f"An unexpected error occurred during text extraction from '{resolved_path}': {str(e)}"
            logger.error(extraction_error, exc_info=True)
            return extraction_error
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF with OCR fallback."""
        try:
            doc = fitz.open(file_path)
            text = "\n".join([page.get_text("text") for page in doc])
            doc.close()
            
            if not text.strip():
                logger.warning("No text found in PDF, attempting OCR...")
                return self._extract_from_scanned_pdf(file_path)
            
            return text
        except Exception as e:
            error_msg = f"Error extracting from PDF: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return error_msg
    
    def _extract_from_scanned_pdf(self, file_path: str) -> str:
        """Extract text from scanned PDF using OCR."""
        try:
            images = convert_from_path(file_path)
            # Aggiunto lang='ita+eng' per migliorare l'accuratezza su testi misti
            text = "\n".join([pytesseract.image_to_string(img, lang='ita+eng') for img in images])
            return text
        except Exception as e:
            error_msg = str(e)
            if "poppler" in error_msg.lower() or "Unable to get page count" in error_msg:
                return (
                    "ERRORE OCR: Per analizzare questo PDF scansionato, è necessario 'Poppler'.\n\n"
                    "Segui questi passaggi per installarlo su Windows:\n"
                    "1. Scarica Poppler da: https://github.com/oschwartz10612/poppler-windows/releases/\n"
                    "2. Estrai l'archivio in una cartella, ad esempio: C:\\Program Files\\poppler\n"
                    "3. Aggiungi il percorso della cartella 'bin' (es. C:\\Program Files\\poppler\\bin) al PATH di sistema.\n"
                    "4. Riavvia il terminale o il tuo IDE."
                )
            return f"Error in PDF OCR: {error_msg}"
    
    def _extract_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(file_path)
            # Aggiunto lang='ita+eng' per migliorare l'accuratezza
            return pytesseract.image_to_string(image, lang='ita+eng')
        except Exception as e:
            return f"Error extracting from image: {str(e)}"
    
    def _extract_from_text(self, file_path: str) -> str:
        """Load text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            return f"Error reading text file from '{file_path}': {str(e)}"
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX (Microsoft Word) files."""
        try:
            logger.info(f"📄 Opening DOCX file: {file_path}")
            doc = Document(file_path)
            
            # Extract text from paragraphs
            full_text = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    full_text.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        full_text.append(" | ".join(row_text))
            
            extracted_text = "\n".join(full_text)
            logger.info(f"📄 Extracted {len(extracted_text)} characters from DOCX")
            
            if not extracted_text.strip():
                return "Warning: No text content found in DOCX file."
            
            return extracted_text
            
        except Exception as e:
            error_msg = f"Error extracting from DOCX: {str(e)}"
            logger.error(error_msg, exc_info=True)
            
            # Provide helpful error for missing python-docx
            if "No module named" in str(e) or "docx" in str(e).lower():
                return (
                    "❌ Errore: La libreria 'python-docx' non è installata.\n\n"
                    "Per installare, esegui:\n"
                    "pip install python-docx\n\n"
                    f"Errore originale: {str(e)}"
                )
            
            return error_msg
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX (Microsoft PowerPoint) files."""
        try:
            from pptx import Presentation
            
            logger.info(f"📊 Opening PPTX file: {file_path}")
            prs = Presentation(file_path)
            
            # Extract text from all slides
            full_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                    
                    # Extract text from tables in shapes
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text)
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                if len(slide_text) > 1:  # More than just the slide header
                    full_text.extend(slide_text)
            
            extracted_text = "\n".join(full_text)
            logger.info(f"📊 Extracted {len(extracted_text)} characters from {len(prs.slides)} slides")
            
            if not extracted_text.strip():
                return "Warning: No text content found in PPTX file."
            
            return extracted_text
            
        except Exception as e:
            error_msg = f"Error extracting from PPTX: {str(e)}"
            logger.error(error_msg, exc_info=True)
            
            # Provide helpful error for missing python-pptx
            if "No module named" in str(e) or "pptx" in str(e).lower():
                return (
                    "❌ Errore: La libreria 'python-pptx' non è installata.\n\n"
                    "Per installare, esegui:\n"
                    "pip install python-pptx\n\n"
                    f"Errore originale: {str(e)}"
                )
            
            return error_msg
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF with OCR fallback."""
        from loguru import logger
        
        try:
            logger.info(f"📄 Opening PDF: {file_path}")
            doc = fitz.open(file_path)
            logger.info(f"📄 PDF opened successfully. Pages: {len(doc)}")
            
            text = "\n".join([page.get_text("text") for page in doc])
            doc.close()
            
            logger.info(f"📄 Extracted {len(text)} characters from PDF")
            
            # Use OCR if no text found (scanned PDF)
            if not text.strip():
                logger.warning("📄 No text found in PDF, attempting OCR...")
                return self._extract_from_scanned_pdf(file_path)
            
            return text
        except Exception as e:
            error_msg = f"Error extracting from PDF: {str(e)}"
            logger.error(error_msg, exc_info=True)
            return error_msg
    
    def _extract_from_scanned_pdf(self, file_path: str) -> str:
        """Extract text from scanned PDF using OCR."""
        from loguru import logger
        
        try:
            logger.info(f"📄 Attempting OCR on PDF: {file_path}")
            images = convert_from_path(file_path)
            logger.info(f"📄 Converted PDF to {len(images)} images")
            
            text = "\n".join([pytesseract.image_to_string(img) for img in images])
            logger.info(f"📄 OCR extracted {len(text)} characters")
            return text
            
        except Exception as e:
            error_msg = str(e)
            logger.error(f"📄 OCR failed: {error_msg}")
            
            # Provide helpful error message for Poppler installation
            if "poppler" in error_msg.lower() or "Unable to get page count" in error_msg:
                return (
                    f"❌ OCR non disponibile: Poppler non è installato.\n\n"
                    f"Il PDF sembra essere scansionato (senza testo estraibile) e richiede OCR.\n\n"
                    f"Per installare Poppler su Windows:\n"
                    f"1. Scarica Poppler da: https://github.com/oschwartz10612/poppler-windows/releases/\n"
                    f"2. Estrai l'archivio in una cartella (es. C:\\poppler)\n"
                    f"3. Aggiungi C:\\poppler\\Library\\bin al PATH di sistema\n\n"
                    f"Errore completo: {error_msg}"
                )
            
            return f"Error in PDF OCR: {error_msg}"
    
    def _extract_from_image(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        try:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        except Exception as e:
            return f"Error extracting from image: {str(e)}"
    
    def _extract_from_text(self, file_path: str) -> str:
        """Load text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            return f"Error reading text file from '{file_path}': {str(e)}\nFile exists: {os.path.exists(file_path)}"


class UniversalDocumentSummarizer(BaseWrapper):
    """Universal document summarization for all file types (PDF, PPTX, DOCX, etc.)."""
    
    def validate_dependencies(self):
        """Validate summarization dependencies."""
        pass
    
    def __init__(self):
        super().__init__()
        self.processor = DocumentProcessor()
    
    def summarize(self, file_path: str) -> str:
        """
        Summarize any document type (PDF, PPTX, DOCX, TXT).
        Extracts text and generates an intelligent summary.
        """
        try:
            from study_buddy.utils.llm import llm
            
            resolved_path = resolve_file_path(file_path)
            logger.info(f"📝 Summarizing document: {resolved_path}")
            
            if not os.path.exists(resolved_path):
                return f"Error: File not found at {resolved_path}"
            
            # Extract text from document
            ext = FileProcessor.get_file_extension(resolved_path)
            logger.info(f"Document type: {ext}")
            
            # Get raw text
            if ext == ".pdf":
                raw_text = self.processor._extract_from_pdf(resolved_path)
            elif ext in [".pptx", ".ppt"]:
                raw_text = self.processor._extract_from_pptx(resolved_path)
            elif ext in [".docx", ".doc"]:
                raw_text = self.processor._extract_from_docx(resolved_path)
            elif ext == ".txt":
                raw_text = self.processor._extract_from_text(resolved_path)
            else:
                return f"Formato non supportato per il riassunto: {ext}"
            
            if raw_text.startswith("Error"):
                return raw_text
            
            # Limit text for summarization (take beginning and key parts)
            text_for_summary = raw_text[:3500] if len(raw_text) > 3500 else raw_text
            
            # Generate summary using LLM
            summary_prompt = f"""Analizza questo documento e crea un riassunto dettagliato in italiano.

DOCUMENTO:
{text_for_summary}

Fornisci un riassunto che includa:
- Di cosa tratta il documento (2-3 righe)
- I punti chiave o sezioni principali
- Eventuali informazioni tecniche rilevanti
- Una breve conclusione

Riassunto:"""

            logger.info("Generating summary with LLM...")
            summary_response = llm.invoke([{"role": "user", "content": summary_prompt}], max_tokens=1500)
            summary_text = summary_response.content if hasattr(summary_response, 'content') else str(summary_response)
            
            # Handle case where content might be a list of dicts
            if isinstance(summary_text, list):
                # Extract 'text' field from dict items if present
                text_parts = []
                for item in summary_text:
                    if isinstance(item, dict) and 'text' in item:
                        text_parts.append(item['text'])
                    else:
                        text_parts.append(str(item))
                summary_text = '\n'.join(text_parts)
            # Handle case where content is a single dict with 'text' key
            elif isinstance(summary_text, dict) and 'text' in summary_text:
                summary_text = summary_text['text']
            summary_text = str(summary_text)  # Ensure it's a string
            
            # Clean the summary from encoded data/signatures
            if summary_text:
                lines = summary_text.split('\n')
                cleaned_lines = []
                for line in lines:
                    # Skip lines that look like encoded data (long strings without spaces)
                    if len(line) > 100 and ' ' not in line[:50]:
                        logger.warning(f"Skipping suspicious encoded line in summary (length: {len(line)})")
                        continue
                    # Skip lines that are just JSON-like metadata
                    if line.strip().startswith('"type":') or line.strip().startswith('{'):
                        continue
                    cleaned_lines.append(line)
                summary_text = '\n'.join(cleaned_lines).strip()
            
            filename = os.path.basename(resolved_path)
            return f"**Riassunto di: {filename}**\n\n{summary_text}"
            
        except Exception as e:
            logger.error(f"Error in universal summarization: {e}", exc_info=True)
            return f"Errore durante il riassunto: {str(e)}"


class DocumentSummarizer(BaseWrapper):
    """Document summarization using Gradio API."""
    
    def validate_dependencies(self):
        """Validate Gradio client dependencies."""
        try:
            from gradio_client import Client, handle_file
        except ImportError:
            raise ValueError("gradio_client not installed. Run: pip install gradio_client")
    
    def __init__(self):
        super().__init__()
        from gradio_client import Client
        self.client = Client("Tulika2000/ai-pdf-summarizer")
    
    def summarize(self, file_path: str, length: str = "Medium", style: str = "Key Takeaways") -> str:
        """
        Summarize document content using Gradio API.
        
        Args:
            file_path: Path to the PDF file
            length: Summary length ("Short", "Medium", "Long")  
            style: Summary style ("Key Takeaways", "Executive Summary", "Detailed Analysis")
        """
        try:
            from gradio_client import handle_file
            
            print(f"Summarizing document: {file_path}")
            resolved_path = resolve_file_path(file_path)
            print(f"Resolved file path: {resolved_path}")
            
            # Verify file exists and is PDF
            if not os.path.exists(resolved_path):
                return f"Error: File not found at {resolved_path}"
            
            ext = FileProcessor.get_file_extension(resolved_path)
            if ext != ".pdf":
                return f"Error: Only PDF files are supported. Got: {ext}"
            
            # Call Gradio API
            result = self.client.predict(
                file=handle_file(resolved_path),
                length=length,
                style=style,
                api_name="/predict"
            )
            
            return result
            
        except Exception as e:
            return f"Error summarizing document: {str(e)}"


class SentimentAnalyzer(BaseWrapper):
    """Text sentiment analysis."""
    
    def validate_dependencies(self):
        """Validate sentiment analysis dependencies."""
        pass
    
    def analyze(self, file_path: str) -> str:
        """Analyze sentiment from file."""
        try:
            resolved_path = resolve_file_path(file_path)
            
            processor = DocumentProcessor()
            text = processor.extract_text(resolved_path)
            
            if text.startswith("Error"):
                return text
            
            blob = TextBlob(text)
            polarity = blob.sentiment.polarity
            subjectivity = blob.sentiment.subjectivity
            
            if polarity > 0:
                sentiment = "Positive"
            elif polarity < 0:
                sentiment = "Negative"
            else:
                sentiment = "Neutral"
            
            return (
                f"Sentiment: {sentiment}\n"
                f"Polarity: {polarity:.3f}\n"
                f"Subjectivity: {subjectivity:.3f}"
            )
            
        except Exception as e:
            return f"Error in sentiment analysis: {str(e)}"


class AudioProcessor(BaseWrapper):
    """Audio processing for TTS and STT."""
    
    def validate_dependencies(self):
        """Validate audio processing dependencies."""
        Config.validate_key("ELEVEN_API_KEY")
        Config.validate_key("ASSEMBLYAI_API_KEY")
    
    def text_to_speech(self, text: str) -> str:
        """Convert text to speech using ElevenLabs."""
        try:
            client = ElevenLabs(api_key=Config.ELEVEN_API_KEY)
            
            audio = client.text_to_speech.convert(
                text=text,
                voice_id=Config.DEFAULT_VOICE_ID,
                model_id=Config.DEFAULT_TTS_MODEL,
                output_format="mp3_44100_128"
            )
            
            temp_path = FileProcessor.create_temp_file(".mp3")
            save(audio, temp_path)
            return temp_path
            
        except Exception as e:
            return f"Error in text-to-speech: {str(e)}"
    
    def speech_to_text(self, audio_input: Union[str, bytes]) -> str:
        """Convert speech to text using AssemblyAI with Italian language support."""
        try:
            import assemblyai as aai
            import os

            # Imposta la chiave API
            aai.settings.api_key = Config.ASSEMBLYAI_API_KEY
            transcriber = aai.Transcriber()

            # Prepara il file audio
            audio_path = self._prepare_audio_path(audio_input)

            # Configurazione della trascrizione
            config = aai.TranscriptionConfig(
                language_code="it",                     
                punctuate=True,                         
                format_text=True,                       
                speech_model=aai.SpeechModel.best       
            )

            # Avvia la trascrizione
            transcript = transcriber.transcribe(audio_path, config)

            # Controlla se c'è stato un errore
            if transcript.status == aai.TranscriptStatus.error:
                raise RuntimeError(f"Transcription failed: {transcript.error}")

            return transcript.text

        except Exception as e:
            return f"Error in speech-to-text: {str(e)}"

        finally:
            # Rimuove il file temporaneo se esiste
            if isinstance(audio_input, (str, bytes)) and os.path.exists(audio_path):
                os.remove(audio_path)

    
    def _prepare_audio_path(self, audio_input: Union[str, bytes]) -> str:
        """Prepare audio file path from various input types."""
        if isinstance(audio_input, str):
            if audio_input.startswith("http"):
                return FileProcessor.download_file(audio_input, ".mp3")
            elif os.path.exists(audio_input):
                return audio_input
            else:
                raise ValueError("Invalid audio path")
        elif isinstance(audio_input, bytes):
            temp_path = FileProcessor.create_temp_file(".mp3")
            with open(temp_path, "wb") as f:
                f.write(audio_input)
            return temp_path
        else:
            raise TypeError("Unsupported audio input type")


class SpotifySearcher(BaseWrapper):
    """Spotify music search functionality."""

    def validate_dependencies(self):
        """Validate Spotify API credentials."""
        Config.validate_key("SPOTIFY_CLIENT_ID")
        Config.validate_key("SPOTIFY_CLIENT_SECRET")

    def __init__(self):
        super().__init__()
        self.access_token = None
        
    def _ensure_token(self):
        if self.access_token is None:
            self.access_token = self._get_access_token()

    def _get_access_token(self) -> str:
        """Get Spotify access token."""
        auth_url = "https://accounts.spotify.com/api/token"
        response = requests.post(
            auth_url,
            data={"grant_type": "client_credentials"},
            auth=(Config.SPOTIFY_CLIENT_ID, Config.SPOTIFY_CLIENT_SECRET),
        )
        response.raise_for_status()
        return response.json()["access_token"]

    def search(self, query: str, search_type: str = "track", limit: int = 5) -> str:
        """Search for music on Spotify."""
        try:
            self._ensure_token()
            url = "https://api.spotify.com/v1/search"
            params = {"q": query, "type": search_type, "limit": limit}
            headers = {"Authorization": f"Bearer {self.access_token}"}

            response = requests.get(url, params=params, headers=headers)
            response.raise_for_status()

            data = response.json()
            items = data.get(f"{search_type}s", {}).get("items", [])

            results = []
            for item in items:
                name = item.get("name", "Unknown")
                artist = item.get("artists", [{}])[0].get("name", "Unknown Artist")
                link = item.get("external_urls", {}).get("spotify", "")
                results.append(f"{name} by {artist} → [Listen]({link})")

            return "\n".join(results) if results else "No results found"

        except Exception as e:
            return f"Error searching Spotify: {str(e)}"


class CodeInterpreter:
    """Code interpreter with E2B sandbox management."""

    def __init__(self):
        self.sandbox = None
    
    def _ensure_sandbox(self):
        """Lazy initialization of E2B sandbox."""
        if self.sandbox is None:
            try:
                self.sandbox = Sandbox.create()
            except Exception as e:
                raise RuntimeError(f"Failed to initialize E2B sandbox: {e}")


    def run_code(self, code: str) -> dict:
        """Execute code in sandbox with basic error handling."""
        try:
            self._ensure_sandbox()
            execution = self.sandbox.run_code(code)
            return {
                "results": getattr(execution, "results", []),
                "stdout": getattr(execution.logs, "stdout", "") if hasattr(execution, "logs") else "",
                "stderr": getattr(execution.logs, "stderr", "") if hasattr(execution, "logs") else "",
                "error": getattr(execution, "error", None),
            }
        except Exception as e:
            raise RuntimeError(f"Code execution failed: {e}")

    def close(self):
        """Clean up sandbox resources."""
        try:
            self.sandbox.kill()
        except:
            pass


class CSVAnalyzer:
    """Hybrid CSV analysis: statistical and semantic analysis."""

    def __init__(self, file_path: str):
        self.file_path = os.path.abspath(file_path)
        self.df = pd.read_csv(self.file_path)
        self.documents = self._load_documents()

    def _load_documents(self) -> List:
        loader = CSVLoader(file_path=self.file_path)
        raw_docs = loader.load()
        splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=50)
        return splitter.split_documents(raw_docs)

    def get_statistical_analysis(self) -> str:
        parts = [
            f"Columns: {list(self.df.columns)}",
            f"Rows: {len(self.df)}",
            f"Data types:\n{self.df.dtypes.to_string()}",
            f"\nFirst 3 rows:\n{self.df.head(3).to_string(index=False)}",
            f"\nDescriptive statistics:\n{self.df.describe(include='all').fillna('').to_string()}",
        ]
        return "\n\n".join(parts)

    def get_semantic_analysis(self) -> str:
        try:
            llm = LLMFactory.create_together_llm()
            content = "\n\n".join([doc.page_content for doc in self.documents[:10]])
            prompt = f"Analyze this CSV dataset content and provide key insights:\n\n{content}"
            return llm.invoke(prompt)
        except Exception as e:
            return f"Error in semantic analysis: {str(e)}"

    def full_analysis(self) -> str:
        return f"[STATISTICAL ANALYSIS]\n{self.get_statistical_analysis()}\n\n[SEMANTIC ANALYSIS]\n{self.get_semantic_analysis()}"


class DataVisualizer:
    """Data visualization tool using sandboxed Python execution."""

    def __init__(self, output_dir: str = "./visualizations"):
        self.output_dir = os.path.abspath(output_dir)
        os.makedirs(self.output_dir, exist_ok=True)
        self.sandbox = None
        self.llm = TogetherClient(api_key=Config.TOGETHER_API_KEY)
        
    def _ensure_sandbox(self):
        if self.sandbox is None:
            self.sandbox = Sandbox.create()

    def create_visualization(self, csv_path: str, query: str) -> dict:
        """Generate visualizations from natural language query."""
        try:
            self._ensure_sandbox()
            resolved_path = os.path.abspath(csv_path)
            with open(resolved_path, "rb") as f:
                sandbox_path = self.sandbox.files.write("dataset.csv", f).path

            code = self._generate_code(query, sandbox_path, resolved_path)
            image_paths = self._execute_and_save(code)
            return {"image_paths": image_paths, "success": True}

        except Exception as e:
            return {"error": str(e), "success": False}

    def _generate_code(self, query: str, sandbox_path: str, local_path: str) -> str:
        """Generate Python code for visualization."""
        df_sample = pd.read_csv(local_path, nrows=300)
        column_info = df_sample.dtypes.astype(str).to_dict()

        prompt = f"""Generate Python code for visualization.
Dataset path in sandbox: {sandbox_path}
Columns and types: {column_info}
User query: "{query}"
Requirements:
- Load data: pd.read_csv("{sandbox_path}")
- Use matplotlib/seaborn
- Save plots: plt.savefig("chart-<name>.png")
- Call plt.close() after saving
- Output ONLY Python code, without ``` and useless spaces
Code:"""

        response = self.llm.chat.completions.create(
            model="meta-llama/Meta-Llama-3.1-8B-Instruct-Turbo",
            messages=[
                {"role": "system", "content": "You are a Python data visualization expert."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1024,
            temperature=0.4
        )

        code = response.choices[0].message.content.strip()
        return self._clean_code(code)

    def _clean_code(self, code: str) -> str:
        """Remove markdown formatting from generated code."""
        if code.startswith("```"):
            lines = code.split("\n")
            if lines[0].startswith("```"):
                lines = lines[1:]
            if lines and lines[-1].strip() == "```":
                lines = lines[:-1]
            code = "\n".join(lines)
        return code.strip()

    def _execute_and_save(self, code: str) -> List[str]:
        """Execute code in sandbox and save images."""
        
        execution = self.sandbox.run_code(code)
        print("Code generated by LLM:", code)
        
        if execution.error:
            raise RuntimeError(f"Code execution failed: {execution.error.value}")

        saved_paths = []
        for result in getattr(execution, "results", []):
            if getattr(result, "png", None):
                filename = f"{uuid.uuid4().hex}.png"
                file_path = os.path.join(self.output_dir, filename)
                with open(file_path, "wb") as f:
                    f.write(base64.b64decode(result.png))
                saved_paths.append(os.path.abspath(file_path))

        if not saved_paths:
            raise RuntimeError("No visualizations were generated")

        return saved_paths

    def close(self):
        """Clean up sandbox resources."""
        try:
            self.sandbox.kill()
        except:
            pass


# =============================================================================
# Custom API Wrappers with Error Handling
# =============================================================================

class GoogleBooksWrapper(GoogleBooksAPIWrapper):
    """Google Books wrapper with better error handling."""
    
    def _format(self, query: str, books: list) -> str:
        """Format search results with graceful handling of missing fields."""
        if not books:
            return f"No books found for query: {query}"
        
        results = [f"Found {len(books)} books for '{query}':"]
        
        for i, book in enumerate(books, start=1):
            info = book.get("volumeInfo", {})
            title = info.get("title", "Title unavailable")
            authors = self._format_authors(info.get("authors", ["Unknown author"]))
            summary = info.get("description", "No description")
            link = info.get("infoLink", "No link available")
            
            result = f'{i}. "{title}" by {authors}\n   {summary}\n   More info: {link}'
            results.append(result)
        
        return "\n\n".join(results)


class GoogleScholarWrapper(GoogleScholarAPIWrapper):
    """Google Scholar wrapper with URL inclusion."""
    
    def run(self, query: str) -> str:
        """Run query with result formatting."""
        total_results = []
        page = 0
        
        while page < max((self.top_k_results - 20), 1):
            results = (
                self.google_scholar_engine({
                    "q": query,
                    "start": page,
                    "hl": self.hl,
                    "num": min(self.top_k_results, 20),
                    "lr": self.lr,
                })
                .get_dict()
                .get("organic_results", [])
            )
            
            total_results.extend(results)
            if not results:
                break
            page += 20
        
        if not total_results:
            return "No Google Scholar results found"
        
        formatted_results = []
        for result in total_results:
            title = result.get('title', 'No title')
            authors = result.get('publication_info', {}).get('authors', [])
            author_names = ', '.join([author.get('name', '') for author in authors])
            summary = result.get('publication_info', {}).get('summary', 'No summary')
            citations = result.get('inline_links', {}).get('cited_by', {}).get('total', 'N/A')
            url = result.get('link', 'No URL available')
            
            formatted_result = (
                f"Title: {title}\n"
                f"Authors: {author_names or 'Unknown'}\n"
                f"Summary: {summary}\n"
                f"Citations: {citations}\n"
                f"URL: {url}"
            )
            formatted_results.append(formatted_result)
        
        return "\n\n".join(formatted_results)


class WikidataSearcher(BaseWrapper):
    """Wikidata SPARQL query wrapper."""
    
    def validate_dependencies(self):
        """Validate Wikidata access."""
        pass
    
    def search(self, query: str) -> str:
        """Search Wikidata using SPARQL."""
        try:
            url = "https://query.wikidata.org/sparql"
            headers = {"User-Agent": "StudyBuddyBot/1.0"}
            
            sparql_query = f"""
            SELECT ?item ?itemLabel ?description ?alias ?propertyLabel ?propertyValueLabel WHERE {{
              ?item ?label "{query}"@en.
              OPTIONAL {{ ?item skos:altLabel ?alias FILTER (LANG(?alias) = "en") }}
              OPTIONAL {{ ?item schema:description ?description FILTER (LANG(?description) = "en") }}
              OPTIONAL {{ 
                ?item ?property ?propertyValue.
                ?property wikibase:directClaim ?propClaim.
                ?propClaim rdfs:label ?propertyLabel.
                ?propertyValue rdfs:label ?propertyValueLabel.
                FILTER(LANG(?propertyLabel) = "en" && LANG(?propertyValueLabel) = "en") 
              }}
              SERVICE wikibase:label {{ bd:serviceParam wikibase:language "en". }}
            }}
            LIMIT 10
            """
            
            response = requests.get(
                url, 
                params={"query": sparql_query, "format": "json"}, 
                headers=headers,
                timeout=10
            )
            response.raise_for_status()
            
            return self._format_results(response.json(), query)
            
        except Exception as e:
            return f"Error searching Wikidata: {str(e)}"
    
    def _format_results(self, data: dict, query: str) -> str:
        """Format Wikidata SPARQL results."""
        results = data.get("results", {}).get("bindings", [])
        
        if not results:
            return f"No Wikidata results for '{query}'"
        
        entities = {}
        for result in results:
            entity_id = result["item"]["value"].split("/")[-1]
            
            if entity_id not in entities:
                entities[entity_id] = {
                    "label": result.get("itemLabel", {}).get("value", "Unknown"),
                    "description": result.get("description", {}).get("value", "No description"),
                    "aliases": set(),
                    "properties": {}
                }
            
            # Collect aliases and properties
            alias = result.get("alias", {}).get("value")
            if alias:
                entities[entity_id]["aliases"].add(alias)
            
            prop_label = result.get("propertyLabel", {}).get("value")
            prop_value = result.get("propertyValueLabel", {}).get("value")
            if prop_label and prop_value:
                if prop_label not in entities[entity_id]["properties"]:
                    entities[entity_id]["properties"][prop_label] = set()
                entities[entity_id]["properties"][prop_label].add(prop_value)
        
        # Format final output
        formatted_results = []
        for entity_id, details in entities.items():
            result_str = [
                f"Entity {entity_id}:",
                f"Label: {details['label']}",
                f"Description: {details['description']}"
            ]
            
            if details["aliases"]:
                result_str.append(f"Aliases: {', '.join(sorted(details['aliases']))}")
            
            for prop, values in details["properties"].items():
                result_str.append(f"{prop}: {', '.join(sorted(values))}")
            
            formatted_results.append("\n".join(result_str))
        
        return "\n\n".join(formatted_results)


# =============================================================================
# Tool Creation Functions
# =============================================================================

def create_basic_tools() -> List[Tool]:
    """Create basic search and information tools."""
    # Initialize API wrappers
    
    # Configure Tavily Search with API key if available
    tavily_kwargs = {
        "max_results": 5,
        "search_depth": "advanced",
        "include_answer": True,
        "include_raw_content": True,
        "include_images": True
    }
    
    if Config.TAVILY_API_KEY:
        tavily_kwargs["tavily_api_key"] = Config.TAVILY_API_KEY
    else:
        logger.warning("TAVILY_API_KEY not set in Config. Web search may fail.")

    web_search = TavilySearch(**tavily_kwargs)
    
    def run_web_search(query: str) -> str:
        """Wrapper for web search to avoid type hint issues in the library."""
        return web_search.run(query)
    
    youtube_tool = YouTubeSearchTool()
    def run_youtube_search(query: str) -> str:
        """Wrapper for YouTube search."""
        return youtube_tool.run(query)

    wikipedia_tool = WikipediaQueryRun(api_wrapper=WikipediaAPIWrapper())
    def run_wikipedia(query: str) -> str:
        """Wrapper for Wikipedia search."""
        return wikipedia_tool.run(query)

    scholar_tool = GoogleScholarQueryRun(api_wrapper=GoogleScholarWrapper())
    def run_scholar(query: str) -> str:
        """Wrapper for Google Scholar search."""
        return scholar_tool.run(query)

    retriever = VectorStoreRetriever()
    wikidata_searcher = WikidataSearcher()
    
    google_lens = None
    try:
        # Prefer GOOGLE_LENS_API_KEY but allow SERP_API_KEY for backward compatibility
        lens_key = None
        if Config.GOOGLE_LENS_API_KEY:
            lens_key = Config.GOOGLE_LENS_API_KEY
        elif Config.SERP_API_KEY:
            lens_key = Config.SERP_API_KEY

        if not lens_key:
            # Will be caught by the except below and a fallback tool will be registered
            raise ValueError("GOOGLE_LENS_API_KEY or SERP_API_KEY must be set in environment variables")

        lens_api_wrapper = GoogleLensAPIWrapper(serp_api_key=lens_key)
        
        def run_google_lens_analysis(query: str) -> str:
            """
            Analyzes an image from a file path or URL using Google Lens.
            
            For local files, automatically uploads them to imgbb (temporary hosting)
            to get a public URL that Google Lens can access. Falls back to OCR if
            upload fails or IMGBB_API_KEY is not set.
            """
            resolved_path = resolve_file_path(query)
            logger.info(f"Analyzing image with Google Lens: {resolved_path}")

            # Check if it's already a public URL
            if resolved_path.startswith("http://") or resolved_path.startswith("https://"):
                logger.info(f"Using public URL directly: {resolved_path}")
                raw_result = lens_api_wrapper.run(resolved_path)
                
                # Parse and structure the Google Lens output for better LLM understanding
                structured_result = "🔍 GOOGLE LENS ANALYSIS - IMAGE SUCCESSFULLY ANALYZED\n\n"
                
                # Try to extract main visual descriptions from the results
                # Look for common patterns in Google Lens results (titles, text matches, etc.)
                lines = raw_result.split('\n')
                visual_subjects = []
                
                # Extract visual information from titles
                for line in lines:
                    if 'Title:' in line:
                        title = line.split('Title:')[1].strip()
                        # Look for common animal/object descriptions
                        keywords = ['dog', 'cat', 'puppy', 'kitten', 'terrier', 'animal', 'bird', 'car', 'building', 'person', 'food']
                        for keyword in keywords:
                            if keyword.lower() in title.lower() and keyword not in visual_subjects:
                                visual_subjects.append(keyword.title())
                
                if visual_subjects:
                    structured_result += f"DETECTED VISUAL CONTENT: {', '.join(visual_subjects)}\n\n"
                
                # Truncate the raw detailed results but keep the summary clear
                if len(raw_result) > 2500:
                    logger.warning(f"Google Lens returned {len(raw_result)} chars, truncating details to 2500")
                    structured_result += "DETAILED SEARCH RESULTS (showing first 2500 chars):\n" + raw_result[:2500] + "\n\n[...results truncated for context limit...]"
                else:
                    structured_result += "DETAILED SEARCH RESULTS:\n" + raw_result
                
                return structured_result
            
            # Check if local file exists
            if not os.path.exists(resolved_path):
                return f"Error: File not found for Google Lens analysis at '{resolved_path}'"
            
            # Verify it's an image file
            ext = FileProcessor.get_file_extension(resolved_path)
            if ext not in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif", ".webp"]:
                return f"Error: File '{os.path.basename(resolved_path)}' is not a supported image format (got {ext})"
            
            # Try to upload local file to imgbb for public access
            public_url = FileProcessor.upload_image_to_imgbb(resolved_path)
            
            if public_url:
                # Successfully uploaded, use Google Lens with public URL
                logger.info(f"Analyzing uploaded image with Google Lens: {public_url}")
                try:
                    raw_result = lens_api_wrapper.run(public_url)
                    
                    # Parse and structure the Google Lens output for better LLM understanding
                    structured_result = "🔍 GOOGLE LENS ANALYSIS - IMAGE SUCCESSFULLY ANALYZED\n\n"
                    
                    # Try to extract main visual descriptions from the results
                    # Look for common patterns in Google Lens results (titles, text matches, etc.)
                    lines = raw_result.split('\n')
                    visual_subjects = []
                    
                    # Extract visual information from titles
                    for line in lines:
                        if 'Title:' in line:
                            title = line.split('Title:')[1].strip()
                            # Look for common animal/object descriptions
                            keywords = ['dog', 'cat', 'puppy', 'kitten', 'terrier', 'animal', 'bird', 'car', 'building', 'person', 'food']
                            for keyword in keywords:
                                if keyword.lower() in title.lower() and keyword not in visual_subjects:
                                    visual_subjects.append(keyword.title())
                    
                    if visual_subjects:
                        structured_result += f"DETECTED VISUAL CONTENT: {', '.join(visual_subjects)}\n\n"
                    
                    # Truncate the raw detailed results but keep the summary clear
                    if len(raw_result) > 2500:
                        logger.warning(f"Google Lens returned {len(raw_result)} chars, truncating details to 2500")
                        structured_result += "DETAILED SEARCH RESULTS (showing first 2500 chars):\n" + raw_result[:2500] + "\n\n[...results truncated for context limit...]"
                    else:
                        structured_result += "DETAILED SEARCH RESULTS:\n" + raw_result
                    
                    return structured_result
                except Exception as e:
                    logger.error(f"Google Lens analysis failed: {e}", exc_info=True)
                    # Fall through to OCR fallback
            else:
                logger.warning("Image upload failed or IMGBB_API_KEY not set, falling back to OCR")
            
            # Fallback: use OCR for local images when upload not available
            try:
                processor = DocumentProcessor()
                ocr_text = processor.extract_text(resolved_path)
                return (
                    "Note: Google Lens requires a public URL. Image upload failed or is not configured "
                    "(set IMGBB_API_KEY to enable automatic uploads). "
                    "Performed OCR on the local image instead:\n\n"
                    f"{ocr_text}"
                )
            except Exception as e2:
                logger.error(f"OCR fallback failed: {e2}", exc_info=True)
                return (
                    f"Error: Could not analyze local image. Upload to temporary hosting failed "
                    f"and OCR fallback encountered an error: {str(e2)}"
                )

        google_lens = StructuredTool.from_function(
            func=run_google_lens_analysis,
            name="google_lens_analyze",
            description="Use this tool to understand and describe the visual content of an image. It identifies objects, scenes, people, places, and landmarks. Use this when the user asks 'what is in this picture?' or 'describe this image'. Do NOT use it to read text from an image.",
            args_schema=GoogleLensInput
        )
    except ValueError as e:
        # If the SERP_API_KEY is not set, register a graceful fallback tool so
        # the system can still handle model requests for 'google_lens_analyze'
        # without failing. The fallback informs the user that Google Lens is
        # not configured and, when possible, performs OCR on local images to
        # return useful information.
        logger.warning(f"Google Lens tool not loaded: {e}")

        def run_google_lens_analysis_stub(query: str) -> str:
            """Fallback for Google Lens analysis when SERP API key is missing.

            Attempts to resolve the provided path. If the file is an image
            present locally, performs OCR (via DocumentProcessor.extract_text)
            and returns the extracted text along with a message indicating
            that full Google Lens features are disabled.
            """
            resolved_path = resolve_file_path(query)
            logger.info(f"Google Lens fallback called with: {resolved_path}")

            # If file exists and is an image, try OCR to provide something useful
            if os.path.exists(resolved_path):
                try:
                    processor = DocumentProcessor()
                    ext = FileProcessor.get_file_extension(resolved_path)
                    if ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif", ".webp"]:
                        ocr_text = processor.extract_text(resolved_path)
                        return (
                            "Google Lens is not configured (GOOGLE_LENS_API_KEY/SERP_API_KEY missing). "
                            "I performed OCR on the provided image and found the following text:\n\n"
                            f"{ocr_text}\n\n"
                            "To enable full Google Lens visual analysis (objects, scenes, landmarks, etc.), "
                            "set GOOGLE_LENS_API_KEY or SERP_API_KEY in your .env file. "
                            "Optionally, set IMGBB_API_KEY to enable automatic image uploads for local files."
                        )
                    else:
                        return (
                            "Google Lens is not configured (GOOGLE_LENS_API_KEY/SERP_API_KEY missing). "
                            f"Received file: {os.path.basename(resolved_path)} (type: {ext}). "
                            "Please set the required API key in your .env file to enable Google Lens analysis."
                        )
                except Exception as e2:
                    logger.error(f"Fallback Google Lens OCR failed: {e2}", exc_info=True)
                    return (
                        "Google Lens is not configured (GOOGLE_LENS_API_KEY/SERP_API_KEY missing). "
                        "Additionally, an internal error occurred while attempting local OCR: "
                        f"{str(e2)}"
                    )

            # File doesn't exist or is a remote URL — inform the user how to enable
            return (
                "Google Lens is not configured (GOOGLE_LENS_API_KEY/SERP_API_KEY missing). "
                "Please set the required API key in your .env file to enable Google Lens analysis. "
                f"Requested resource: {query}"
            )

        google_lens = StructuredTool.from_function(
            func=run_google_lens_analysis_stub,
            name="google_lens_analyze",
            description=(
                "Fallback Google Lens tool: Google Lens is not configured because "
                "SERP_API_KEY is missing. This stub will attempt OCR on local images "
                "and otherwise return instructions to enable the full Google Lens integration."
            ),
            args_schema=GoogleLensInput
        )
    
    tools = [
        Tool(
            name="retrieve_knowledge",
            description="PRIMARY TOOL for answering questions. Search the local knowledge base (slides, syllabus, docs) for information about courses, professors, exams, etc. ALWAYS use this first for any question not about a specific file path.",
            func=retriever.retrieve # Assumendo che retrieve restituisca una stringa
        ),
        Tool(
            name="web_search",
            description="Search the web for current information, news, and real-time updates.",
            func=run_web_search
        ),
        Tool(
            name="youtube_search",
            description="Search for educational videos on YouTube",
            func=run_youtube_search
        ),
        Tool(
            name="wikipedia_search",
            description="Search Wikipedia for encyclopedic information",
            func=run_wikipedia
        ),
        Tool(
            name="google_scholar_search",
            description="Search academic literature on Google Scholar",
            func=run_scholar
        ),
        Tool(
            name="wikidata_search",
            description="Search structured data on Wikidata",
            func=wikidata_searcher.search
        )
    ]
    if google_lens:
        tools.append(google_lens)
    
    return tools


def create_document_tools() -> List[Tool]:
    """Create document processing tools."""
    processor = DocumentProcessor()
    universal_summarizer = UniversalDocumentSummarizer()
    sentiment_analyzer = SentimentAnalyzer()
    
    return [
        Tool(
            name="summarize_document",
            description="Riassumi qualsiasi documento (PDF, PPTX, DOCX, TXT). Usa questo tool quando l'utente chiede un riassunto. NON usare extract_text per riassunti - usa direttamente questo tool.",
            func=universal_summarizer.summarize
        ),
        StructuredTool.from_function(
            func=processor.extract_text,
            name="extract_text",
            description="Extract RAW TEXT from a specific file path ONLY when the user explicitly asks to READ/VIEW a file by name. DO NOT use for general questions -> use retrieve_knowledge instead.",
            args_schema=FilePathInput
        ),
        Tool(
            name="analyze_sentiment",
            description="Analyze sentiment, polarity, and subjectivity of document content",
            func=sentiment_analyzer.analyze
        )
    ]


def create_audio_tools() -> List[Tool]:
    """Create audio processing tools."""
    try:
        audio_processor = AudioProcessor()
        
        return [
            Tool(
                name="text_to_speech",
                description="Convert text to speech using ElevenLabs TTS",
                func=audio_processor.text_to_speech
            ),
            Tool(
                name="speech_to_text",
                description="Transcribe audio files to text using AssemblyAI",
                func=audio_processor.speech_to_text
            )
        ]
    except ValueError as e:
        print(f"Warning: Audio tools not available - {e}")
        return []


def create_multimedia_tools() -> List[Tool]:
    """Create multimedia and entertainment tools."""
    tools = []
    
    # Spotify search
    try:
        spotify = SpotifySearcher()
        tools.append(Tool(
            name="spotify_search",
            description="Search for music on Spotify and get listening links",
            func=spotify.search
        ))
    except ValueError as e:
        print(f"Warning: Spotify search not available - {e}")
    
    return tools



def create_data_analysis_tools() -> List[Tool]:
    """Create data analysis and visualization tools."""
    tools = []

    # Code interpreter 
    try:
        code_interpreter = CodeInterpreter()
        
        def execute_code(code: str) -> str:
            """Safely execute code and return formatted results."""
            try:
                result = code_interpreter.run_code(code)
                
                # Format the execution results
                output_parts = []
                
                if result.get("stdout"):
                    output_parts.append(f"Output:\n{result['stdout']}")
                
                if result.get("stderr"):
                    output_parts.append(f"Errors:\n{result['stderr']}")
                
                if result.get("error"):
                    error_msg = result['error']
                    if hasattr(error_msg, 'message'):
                        error_text = error_msg.message
                    elif hasattr(error_msg, 'value'):
                        error_text = error_msg.value
                    else:
                        error_text = str(error_msg)
                    output_parts.append(f"Execution Error:\n{error_text}")
                
                if result.get("results"):
                    # Handle different types of results (charts, data, etc.)
                    for i, res in enumerate(result["results"]):
                        if hasattr(res, 'png') and res.png:
                            output_parts.append(f"Generated visualization {i+1} (PNG image)")
                        elif hasattr(res, 'text') and res.text:
                            output_parts.append(f"Text output {i+1}:\n{res.text}")
                        elif hasattr(res, 'json') and res.json:
                            output_parts.append(f"JSON output {i+1}:\n{res.json}")
                        else:
                            output_parts.append(f"Result {i+1}: {str(res)}")
                
                return "\n\n".join(output_parts) if output_parts else "Code executed successfully with no output"
                
            except Exception as e:
                return f"Error executing code: {str(e)}"
        
        tools.append(StructuredTool.from_function(
            func=execute_code,
            name="execute_code",
            description="Execute Python code in a secure E2B sandbox environment. Supports data analysis, visualization, file operations, and more.",
            args_schema=CodeInterpreterInput
        ))
        
    except ValueError as e:
        print(f"Warning: Code interpreter not available - {e}")
    except Exception as e:
        print(f"Warning: Failed to initialize code interpreter - {e}")

    # CSV analysis tool
    def analyze_csv(file_path: str) -> str:
        """Safely analyze CSV files."""
        try:
            resolved_path = resolve_file_path(file_path)
            if not os.path.exists(resolved_path):
                return f"Error: File not found at {resolved_path}"
            
            analyzer = CSVAnalyzer(resolved_path)
            return analyzer.full_analysis()
        except Exception as e:
            return f"Error analyzing CSV: {str(e)}"

    tools.append(StructuredTool.from_function(
        func=analyze_csv,
        name="analyze_csv",
        description="Perform statistical and semantic analysis on CSV files. Provides data overview, statistics, and AI-generated insights.",
        args_schema=FilePathInput
    ))

    # Data visualization tool
    try:
        visualizer = DataVisualizer()

        def create_visualization(csv_path: str, query: str) -> str:
            """Safely create data visualizations."""
            try:
                print(f"Creating visualization for file: {csv_path} with query: {query}")
                resolved_path = resolve_file_path(csv_path)
                print(f"Resolved file path: {resolved_path}")
                
                if not os.path.exists(resolved_path):
                    return f"Error: The file '{resolved_path}' doesn't exist. Please upload a valid CSV file."
                
                result = visualizer.create_visualization(resolved_path, query)
                
                if result.get("success"):
                    image_paths = result.get("image_paths", [])
                    if image_paths:
                        paths_str = "\n".join([f"- {path}" for path in image_paths])
                        return f"Successfully created {len(image_paths)} visualization(s):\n{paths_str}"
                    else:
                        return "Visualization completed but no image files were generated."
                else:
                    return f"Visualization failed: {result.get('error', 'Unknown error')}"
                    
            except Exception as e:
                return f"Error creating visualization: {str(e)}"

        tools.append(StructuredTool.from_function(
            func=create_visualization,
            name="create_visualization",
            description="Generate data visualizations from CSV files using natural language queries. Creates charts, graphs, and plots based on your description.",
            args_schema=DataVizInput
        ))

    except ValueError as e:
        print(f"Warning: Data visualizer not available - {e}")
    except Exception as e:
        print(f"Warning: Failed to initialize data visualizer - {e}")

    return tools



# =============================================================================
# Main Tool Assembly
# =============================================================================

def get_all_tools() -> List[Tool]:
    """Assemble all available tools with enhanced logging and error handling."""
    all_tools = []
    
    with LogContext("tool_initialization", logger) as log_ctx:
        logger.info("🔧 Starting tool initialization")
        
        tool_categories = {
            "basic": create_basic_tools,
            "document": create_document_tools,
            "audio": create_audio_tools,
            "multimedia": create_multimedia_tools,
            "data": create_data_analysis_tools
        }
        
        for category, creator_func in tool_categories.items():
            try:
                with LogContext(f"{category}_tools", logger):
                    tools = creator_func()
                    all_tools.extend(tools)
                    logger.success(f"✅ Loaded {len(tools)} {category} tools")
                    
                    # Log individual tool details
                    for tool in tools:
                        logger.debug(f"  - {tool.name}: {tool.description[:100]}...")
                        
            except Exception as e:
                logger.error(f"❌ Failed to load {category} tools: {str(e)}")
                logger.exception(e)  # Log full traceback
        
        # Log final tool initialization summary
        logger.success(f"""
📊 Tool Initialization Summary:
   Total Tools: {len(all_tools)}
   Categories Loaded: {sum(1 for t in tool_categories if any(tool for tool in all_tools if tool.name.startswith(t)))}
   Tool Names: {[tool.name for tool in all_tools]}
""")
    
    return all_tools
